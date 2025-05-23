import logging
import os
import sys
import traceback
from typing import Any, Dict, List, Optional, Union

from dotenv import load_dotenv
from mcp.server.fastmcp import FastMCP
from pydantic import Field

from aworld.logs.util import logger

# PPTX processing library
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
    from pptx.shapes.autoshape import Shape as PptxShape  # pylint: disable=W0611
    from pptx.table import Table as PptxTable
except ImportError:
    logging.error(
        "python-pptx library is not installed. "
        "PPTX operations will not work. Please install it by running: pip install python-pptx"
    )
    Presentation = None
    PP_PLACEHOLDER_TYPE = None
    MSO_SHAPE_TYPE = None
    PptxShape = None
    PptxTable = None

# pylint: disable=W0707
# Initialize MCP server
mcp = FastMCP("pptx-server")


def _get_slide_indices(prs: Any, page_numbers: Optional[List[int]]) -> List[int]:
    """Helper to get valid slide indices to process."""
    if Presentation is None:  # Should not happen if initial check passes
        return []
    total_slides = len(prs.slides)
    if page_numbers:
        valid_indices = []
        for pn in page_numbers:
            if 0 <= pn < total_slides:
                valid_indices.append(pn)
            else:
                logger.warning(f"Page number {pn} is out of range for PPTX with {total_slides} slides. Skipping.")
        return valid_indices
    return list(range(total_slides))


@mcp.tool(
    description=(
        "Extracts text content from a PPTX file. "
        "Supports specific slide numbers, character limits, and page-by-page output."
    )
)
async def extract_pptx_text(
    file_path: str = Field(description="The absolute path to the PPTX file."),
    slide_numbers: Optional[List[int]] = Field(
        default=None,
        description="Optional. A list of slide numbers (0-indexed) to extract text from. "
        "If not provided, extracts from all slides.",
    ),
    max_chars: Optional[int] = Field(
        default=None,
        description="Optional. Maximum number of characters to return. "
        "Text will be truncated if it exceeds this limit.",
    ),
    return_format: str = Field(
        default="full_text",
        description="Optional. 'full_text' (default) to return a single string, "
        "or 'by_slide' to return a dictionary of slide_number:text.",
        enum=["full_text", "by_slide"],
    ),
    include_speaker_notes: bool = Field(
        default=False, description="Optional. Whether to include speaker notes in the extracted text."
    ),
) -> Union[str, Dict[str, str]]:
    """
    Extracts text content from a PPTX file, including text from shapes, tables, and optionally speaker notes.

    Args:
        file_path: The absolute path to the PPTX file.
        slide_numbers: Optional. A list of 0-indexed slide numbers to extract text from.
                       If not provided, extracts from all slides.
        max_chars: Optional. Maximum number of characters to return. Text will be truncated.
        return_format: 'full_text' for a single string, 'by_slide' for a dict of slide_number:text.
        include_speaker_notes: If True, includes speaker notes in the output.

    Returns:
        A string containing the full text or a dictionary mapping slide numbers (e.g., "slide_1")
        to text content, depending on `return_format`.

    Raises:
        RuntimeError: If python-pptx is not installed.
        FileNotFoundError: If the specified file_path does not exist.
        Exception: If an error occurs during processing.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx library is required for text extraction but not found.")
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        prs = Presentation(file_path)
        slides_to_process_indices = _get_slide_indices(prs, slide_numbers)
        current_char_count = 0

        if return_format == "by_slide":
            slide_text_map = {}
            for slide_idx in slides_to_process_indices:
                slide = prs.slides[slide_idx]
                slide_text_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text_parts.append(shape.text_frame.text)
                    elif shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                slide_text_parts.append(cell.text_frame.text)

                if include_speaker_notes and slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        slide_text_parts.append(f"\n[Speaker Notes]:\n{notes_slide.notes_text_frame.text}")

                text = "\n".join(filter(None, slide_text_parts))

                if max_chars is not None:
                    remaining_chars_for_slide = max_chars - current_char_count
                    if remaining_chars_for_slide <= 0 and len(slides_to_process_indices) > 1:
                        logger.warning(
                            f"Max characters ({max_chars}) reached. Skipping further slides for 'by_slide' format."
                        )
                        break
                    text = text[:remaining_chars_for_slide]
                    current_char_count += len(text)

                slide_text_map[f"slide_{slide_idx + 1}"] = text
                if max_chars is not None and current_char_count >= max_chars and len(slides_to_process_indices) > 1:
                    break
            return slide_text_map
        else:  # full_text
            full_text_content_parts = []
            for slide_idx in slides_to_process_indices:
                slide = prs.slides[slide_idx]
                slide_text_parts = []
                if len(slides_to_process_indices) > 1:
                    full_text_content_parts.append(f"\n--- Slide {slide_idx + 1} ---\n")

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text_parts.append(shape.text_frame.text)
                    elif shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                slide_text_parts.append(cell.text_frame.text)

                if include_speaker_notes and slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        slide_text_parts.append(f"\n[Speaker Notes]:\n{notes_slide.notes_text_frame.text}")

                text_for_this_slide = "\n".join(filter(None, slide_text_parts))

                if max_chars is not None:
                    chars_to_add = len(text_for_this_slide)
                    if current_char_count + chars_to_add > max_chars:
                        chars_to_add = max_chars - current_char_count
                        text_for_this_slide = text_for_this_slide[:chars_to_add]

                    full_text_content_parts.append(text_for_this_slide)
                    current_char_count += chars_to_add
                    if current_char_count >= max_chars:
                        break
                else:
                    full_text_content_parts.append(text_for_this_slide)

            return "".join(full_text_content_parts)

    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Error processing PPTX for text extraction: {str(e)}")


@mcp.tool(description="Retrieves metadata from a PPTX file (e.g., author, title, number of slides).")
async def get_pptx_metadata(
    file_path: str = Field(description="The absolute path to the PPTX file."),
) -> Dict[str, Any]:
    """
    Retrieves metadata from a PPTX file.

    Args:
        file_path: The absolute path to the PPTX file.

    Returns:
        A dictionary containing metadata such as author, title, subject, keywords,
        last_modified_by, revision, and number_of_slides.

    Raises:
        RuntimeError: If python-pptx is not installed.
        FileNotFoundError: If the specified file_path does not exist.
        Exception: If an error occurs during processing.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx library is required for metadata retrieval but not found.")
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        prs = Presentation(file_path)
        core_props = prs.core_properties
        return {
            "author": core_props.author,
            "category": core_props.category,
            "comments": core_props.comments,
            "content_status": core_props.content_status,
            "created": str(core_props.created) if core_props.created else None,
            "identifier": core_props.identifier,
            "keywords": core_props.keywords,
            "language": core_props.language,
            "last_modified_by": core_props.last_modified_by,
            "last_printed": str(core_props.last_printed) if core_props.last_printed else None,
            "modified": str(core_props.modified) if core_props.modified else None,
            "revision": core_props.revision,
            "subject": core_props.subject,
            "title": core_props.title,
            "version": core_props.version,
            "number_of_slides": len(prs.slides),
        }
    except Exception as e:
        logger.error(f"Error getting metadata from {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Error processing PPTX for metadata: {str(e)}")


@mcp.tool(description="Extracts all images from a PPTX file and saves them to a specified directory.")
async def extract_images_from_pptx(
    file_path: str = Field(description="Absolute path to the input PPTX file."),
    output_directory: str = Field(description="Directory to save the extracted images."),
    slide_numbers: Optional[List[int]] = Field(
        default=None,
        description="Optional. A list of 0-indexed slide numbers to extract images from. "
        "If not provided, extracts from all slides.",
    ),
) -> List[str]:
    """
    Extracts all images from specified slides of a PPTX file and saves them.

    Args:
        file_path: Absolute path to the input PPTX file.
        output_directory: Directory to save the extracted images.
        slide_numbers: Optional. List of 0-indexed slide numbers. Extracts from all if None.

    Returns:
        A list of absolute paths to the extracted image files.

    Raises:
        RuntimeError: If python-pptx is not installed.
        FileNotFoundError: If the input file does not exist.
        ValueError: If the output directory cannot be created.
        Exception: If an error occurs during image extraction.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx library is required for image extraction but not found.")
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
    if not os.path.isdir(output_directory):
        try:
            os.makedirs(output_directory, exist_ok=True)
        except Exception as e:
            raise ValueError(f"Cannot create output directory: {output_directory}, error: {str(e)}")

    extracted_image_paths = []
    try:
        prs = Presentation(file_path)
        slides_to_process_indices = _get_slide_indices(prs, slide_numbers)

        for slide_idx in slides_to_process_indices:
            slide = prs.slides[slide_idx]
            img_counter = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # type: ignore
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_filename = os.path.join(
                        output_directory, f"slide{slide_idx + 1}_img{img_counter + 1}.{image_ext}"
                    )
                    with open(image_filename, "wb") as img_file:
                        img_file.write(image_bytes)
                    extracted_image_paths.append(image_filename)
                    img_counter += 1
                # Could also check for MSO_SHAPE_TYPE.GROUP and recurse if needed
        return extracted_image_paths
    except Exception as e:
        logger.error(f"Error extracting images from {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Error extracting images: {str(e)}")


@mcp.tool(description="Gets the total number of slides in a PPTX file.")
async def get_slide_count(file_path: str = Field(description="The absolute path to the PPTX file.")) -> int:
    """
    Returns the total number of slides in a PPTX presentation.

    Args:
        file_path: The absolute path to the PPTX file.

    Returns:
        The total number of slides.

    Raises:
        RuntimeError: If python-pptx is not installed.
        FileNotFoundError: If the specified file_path does not exist.
        Exception: If an error occurs during processing.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx library is required for getting slide count but not found.")
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")
    try:
        prs = Presentation(file_path)
        return len(prs.slides)
    except Exception as e:
        logger.error(f"Error getting slide count for {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Error processing PPTX for slide count: {str(e)}")


@mcp.tool(description="Extracts speaker notes from a PPTX file.")
async def extract_speaker_notes(
    file_path: str = Field(description="The absolute path to the PPTX file."),
    slide_numbers: Optional[List[int]] = Field(
        default=None,
        description="Optional. A list of 0-indexed slide numbers to extract notes from. "
        "If not provided, extracts from all slides with notes.",
    ),
    return_format: str = Field(
        default="full_text",
        description="Optional. 'full_text' (default) to return a single string of all notes, "
        "or 'by_slide' to return a dictionary of slide_number:notes_text.",
        enum=["full_text", "by_slide"],
    ),
    max_chars: Optional[int] = Field(
        default=None,
        description="Optional. Maximum number of characters to return for the entire output. "
        "Notes will be truncated if the total exceeds this limit.",
    ),
) -> Union[str, Dict[str, str]]:
    """
    Extracts speaker notes from specified slides or all slides in a PPTX file.

    Args:
        file_path: The absolute path to the PPTX file.
        slide_numbers: Optional. List of 0-indexed slide numbers. Extracts from all if None.
        return_format: 'full_text' for a single string, 'by_slide' for a dict.
        max_chars: Optional. Maximum total characters to return.

    Returns:
        A string containing all speaker notes or a dictionary mapping slide numbers
        (e.g., "slide_1_notes") to their speaker notes.

    Raises:
        RuntimeError: If python-pptx is not installed.
        FileNotFoundError: If the specified file_path does not exist.
        Exception: If an error occurs during processing.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx library is required for speaker notes extraction but not found.")
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        prs = Presentation(file_path)
        slides_to_process_indices = _get_slide_indices(prs, slide_numbers)
        current_char_count = 0

        if return_format == "by_slide":
            slide_notes_map = {}
            for slide_idx in slides_to_process_indices:
                slide = prs.slides[slide_idx]
                notes_text = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text or ""

                if max_chars is not None:
                    remaining_chars_for_note = max_chars - current_char_count
                    if remaining_chars_for_note <= 0 and len(slides_to_process_indices) > 1:
                        logger.warning(
                            f"Max characters ({max_chars}) reached for speaker notes. "
                            f"Skipping further slides for 'by_slide' format."
                        )
                        break
                    notes_text = notes_text[:remaining_chars_for_note]
                    current_char_count += len(notes_text)

                if notes_text:  # Only add if there are notes
                    slide_notes_map[f"slide_{slide_idx + 1}_notes"] = notes_text

                if max_chars is not None and current_char_count >= max_chars and len(slides_to_process_indices) > 1:
                    break
            return slide_notes_map
        else:  # full_text
            all_notes_parts = []
            for slide_idx in slides_to_process_indices:
                slide = prs.slides[slide_idx]
                notes_text_for_slide = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_text_for_slide = notes_slide.notes_text_frame.text or ""

                if notes_text_for_slide:
                    if len(slides_to_process_indices) > 1 or slide_numbers is not None:
                        header = f"\n--- Speaker Notes for Slide {slide_idx + 1} ---\n"
                    else:
                        header = ""

                    current_note_content = header + notes_text_for_slide

                    if max_chars is not None:
                        chars_to_add = len(current_note_content)
                        if current_char_count + chars_to_add > max_chars:
                            chars_to_add = max_chars - current_char_count
                            # Need to be careful not to cut off header if possible
                            if len(header) < chars_to_add:
                                current_note_content = header + notes_text_for_slide[: chars_to_add - len(header)]
                            else:
                                current_note_content = header[:chars_to_add]

                        all_notes_parts.append(current_note_content)
                        current_char_count += len(current_note_content)
                        if current_char_count >= max_chars:
                            break
                    else:
                        all_notes_parts.append(current_note_content)
            return "".join(all_notes_parts)

    except Exception as e:
        logger.error(f"Error extracting speaker notes from {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Error processing PPTX for speaker notes extraction: {str(e)}")


@mcp.tool(description="Extracts titles of specified slides or all slides from a PPTX file.")
async def get_slide_titles(
    file_path: str = Field(description="The absolute path to the PPTX file."),
    slide_numbers: Optional[List[int]] = Field(
        default=None,
        description="Optional. A list of 0-indexed slide numbers to extract titles from. "
        "If not provided, extracts from all slides.",
    ),
    return_format: str = Field(
        default="list",
        description="Optional. 'list' (default) to return a list of titles, "
        "or 'by_slide' to return a dictionary of slide_number:title.",
        enum=["list", "by_slide"],
    ),
) -> Union[List[str], Dict[str, str]]:
    """
    Extracts titles from specified slides or all slides in a PPTX file.
    A slide title is typically the text in the main title placeholder.

    Args:
        file_path: The absolute path to the PPTX file.
        slide_numbers: Optional. List of 0-indexed slide numbers. Extracts from all if None.
        return_format: 'list' for a list of titles, 'by_slide' for a dict of slide_number:title.

    Returns:
        A list of slide titles or a dictionary mapping slide numbers (e.g., "slide_1_title")
        to their titles.

    Raises:
        RuntimeError: If python-pptx is not installed.
        FileNotFoundError: If the specified file_path does not exist.
        Exception: If an error occurs during processing.
    """
    if Presentation is None or PP_PLACEHOLDER_TYPE is None:
        raise RuntimeError("python-pptx library is required for title extraction but not found.")
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        prs = Presentation(file_path)
        slides_to_process_indices = _get_slide_indices(prs, slide_numbers)

        titles_list = []
        titles_map = {}

        for slide_idx in slides_to_process_indices:
            slide = prs.slides[slide_idx]
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text
            else:  # Fallback: check for common title placeholder types
                for shape in slide.placeholders:
                    if shape.placeholder_format.type in (
                        PP_PLACEHOLDER_TYPE.TITLE,  # type: ignore
                        PP_PLACEHOLDER_TYPE.CENTER_TITLE,  # type: ignore
                        PP_PLACEHOLDER_TYPE.SUBTITLE,  # type: ignore
                        PP_PLACEHOLDER_TYPE.VERTICAL_TITLE,  # type: ignore
                    ):
                        if shape.has_text_frame and shape.text_frame.text:
                            title_text = shape.text_frame.text
                            break

            if return_format == "list":
                titles_list.append(title_text)
            else:  # by_slide
                titles_map[f"slide_{slide_idx + 1}_title"] = title_text

        return titles_map if return_format == "by_slide" else titles_list

    except Exception as e:
        logger.error(f"Error extracting slide titles from {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Error processing PPTX for slide title extraction: {str(e)}")


@mcp.tool(description="Retrieves detailed information about shapes on specified slides of a PPTX file.")
async def get_slide_shapes_details(
    file_path: str = Field(description="The absolute path to the PPTX file."),
    slide_numbers: Optional[List[int]] = Field(
        default=None,
        description="Optional. A list of 0-indexed slide numbers to analyze. If not provided, analyzes all slides.",
    ),
) -> Dict[str, List[Dict[str, Any]]]:
    """
    Retrieves detailed information about each shape on specified slides in a PPTX file.

    Args:
        file_path: The absolute path to the PPTX file.
        slide_numbers: Optional. List of 0-indexed slide numbers. Analyzes all if None.

    Returns:
        A dictionary where keys are slide numbers (e.g., "slide_1_shapes") and values are lists
        of dictionaries, each describing a shape (type, text, geometry, etc.).

    Raises:
        RuntimeError: If python-pptx is not installed.
        FileNotFoundError: If the specified file_path does not exist.
        Exception: If an error occurs during processing.
    """
    if Presentation is None or MSO_SHAPE_TYPE is None:
        raise RuntimeError("python-pptx library is required for shape analysis but not found.")
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        prs = Presentation(file_path)
        slides_to_process_indices = _get_slide_indices(prs, slide_numbers)
        shapes_details_map = {}

        for slide_idx in slides_to_process_indices:
            slide = prs.slides[slide_idx]
            slide_key = f"slide_{slide_idx + 1}_shapes"
            shapes_details_map[slide_key] = []
            for shape_idx, shape in enumerate(slide.shapes):
                shape_info = {
                    "shape_index": shape_idx,
                    "name": shape.name,
                    "type": shape.shape_type.name if shape.shape_type else "UNKNOWN",  # MSO_SHAPE_TYPE
                    "has_text_frame": shape.has_text_frame,
                    "text": shape.text_frame.text if shape.has_text_frame else None,
                    "has_table": shape.has_table,
                    "has_chart": shape.has_chart,
                    "is_placeholder": shape.is_placeholder,
                    "placeholder_type": shape.placeholder_format.type.name
                    if shape.is_placeholder and shape.placeholder_format.type
                    else None,
                    "left_inches": shape.left.inches if shape.left else None,
                    "top_inches": shape.top.inches if shape.top else None,
                    "width_inches": shape.width.inches if shape.width else None,
                    "height_inches": shape.height.inches if shape.height else None,
                }
                shapes_details_map[slide_key].append(shape_info)
        return shapes_details_map

    except Exception as e:
        logger.error(f"Error extracting shape details from {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Error processing PPTX for shape details: {str(e)}")


@mcp.tool(description="Extracts data from tables on specified slides of a PPTX file.")
async def extract_table_data_from_slides(
    file_path: str = Field(description="The absolute path to the PPTX file."),
    slide_numbers: Optional[List[int]] = Field(
        default=None,
        description="Optional. A list of 0-indexed slide numbers to extract tables from. "
        "If not provided, extracts from all slides.",
    ),
    return_format: str = Field(
        default="list_of_lists",
        description="Optional. 'list_of_lists' to return table data as a list of lists (rows of cells), "
        "or 'list_of_dicts' to return table data as a list of dictionaries (rows with header keys).",
        enum=["list_of_lists", "list_of_dicts"],
    ),
) -> Dict[str, List[Union[List[List[str]], List[Dict[str, str]]]]]:
    """
    Extracts data from all tables found on specified slides in a PPTX file.

    Args:
        file_path: The absolute path to the PPTX file.
        slide_numbers: Optional. List of 0-indexed slide numbers. Extracts from all if None.
        return_format: 'list_of_lists' (default) or 'list_of_dicts'. For 'list_of_dicts',
                       the first row of the table is assumed to be the header.

    Returns:
        A dictionary where keys are slide numbers (e.g., "slide_1_tables") and values are lists
        of tables found on that slide. Each table is represented as specified by return_format.

    Raises:
        RuntimeError: If python-pptx is not installed.
        FileNotFoundError: If the specified file_path does not exist.
        Exception: If an error occurs during processing.
    """
    if Presentation is None or PptxTable is None:
        raise RuntimeError("python-pptx library is required for table extraction but not found.")
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        prs = Presentation(file_path)
        slides_to_process_indices = _get_slide_indices(prs, slide_numbers)
        slide_tables_map = {}

        for slide_idx in slides_to_process_indices:
            slide = prs.slides[slide_idx]
            slide_key = f"slide_{slide_idx + 1}_tables"
            slide_tables_map[slide_key] = []
            table_idx_on_slide = 0
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    if return_format == "list_of_dicts":
                        if not table.rows:  # Empty table
                            slide_tables_map[slide_key].append([])
                            continue
                        header_cells = [cell.text_frame.text for cell in table.rows[0].cells]
                        for row_idx in range(1, len(table.rows)):
                            row_dict = {}
                            for col_idx, cell in enumerate(table.rows[row_idx].cells):
                                if col_idx < len(header_cells):
                                    row_dict[header_cells[col_idx]] = cell.text_frame.text
                                else:
                                    row_dict[f"column_{col_idx + 1}"] = (
                                        cell.text_frame.text
                                    )  # Fallback if more cols than headers
                            table_data.append(row_dict)
                    else:  # list_of_lists
                        for row in table.rows:
                            row_data = [cell.text_frame.text for cell in row.cells]
                            table_data.append(row_data)

                    slide_tables_map[slide_key].append(table_data)
                    table_idx_on_slide += 1
        return slide_tables_map

    except Exception as e:
        logger.error(f"Error extracting table data from {file_path}: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Error processing PPTX for table data extraction: {str(e)}")


def main():
    """
    Main function to start the MCP server.
    """
    load_dotenv()  # Load environment variables
    if Presentation is None:
        logger.critical(
            "python-pptx library is not installed. PPTX MCP server cannot start. "
            "Please install it: pip install python-pptx"
        )
        sys.exit(1)
    mcp.run(transport="stdio")


# Make the module callable for uvx
def __call__():
    """
    Make the module callable for uvx.
    This function is called when the module is executed directly.
    """
    main()


# Add this for compatibility with uvx
if __name__ != "__main__":  # Ensure this is also set when imported as a module
    sys.modules[__name__].__call__ = __call__

# Run the server when the script is executed directly
if __name__ == "__main__":
    main()
