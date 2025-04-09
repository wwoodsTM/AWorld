import io
import json
import os
import tempfile
from datetime import date, datetime

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pydantic import Field

from aworld.logs.util import logger
from aworld.mcp_servers.image import encode_image_from_file
from aworld.mcp_servers.utils import run_mcp_server
from aworld.utils import import_package, import_packages

import_package("cv2", install_name="opencv-python")
import_packages(["xmltodict", "pandas", "docx2markdown", "PyPDF2"])


class ComplexEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, datetime):
            return obj.strftime("%Y-%m-%d %H:%M:%S")
        elif isinstance(obj, date):
            return obj.strftime("%Y-%m-%d")
        else:
            return json.JSONEncoder.default(self, obj)


def handle_error(e: Exception, error_type: str) -> str:
    """Unified error handling and return standard format error message"""
    error_msg = f"{error_type} error: {str(e)}"
    logger.error(error_msg)
    return json.dumps({"error": error_msg})


def check_file_readable(document_path: str) -> str:
    """Check if file exists and is readable, return error message or None"""
    if not os.path.exists(document_path):
        return f"File does not exist: {document_path}"
    if not os.access(document_path, os.R_OK):
        return f"File is not readable: {document_path}"
    return None


def mcpreadtext(
    document_path: str = Field(description="The input text file path."),
) -> str:
    """Read and return content from text file."""
    error = check_file_readable(document_path)
    if error:
        return json.dumps({"error": error})

    try:
        with open(document_path, "r", encoding="utf-8") as f:
            content = f.read()
        return json.dumps({"content": content})
    except Exception as e:
        return handle_error(e, "Text file reading")


def mcpreadjson(
    document_path: str = Field(description="Path to JSON or JSONL file"),
    is_jsonl: bool = Field(
        default=False,
        description="Whether the file is in JSONL format (one JSON object per line)",
    ),
) -> str:
    """Read and parse JSON or JSONL file, return the parsed content."""
    error = check_file_readable(document_path)
    if error:
        return json.dumps({"error": error})

    try:
        # Choose processing method based on file type
        if is_jsonl:
            # Process JSONL file (one JSON object per line)
            results = []
            with open(document_path, "r", encoding="utf-8") as f:
                for line_num, line in enumerate(f, 1):
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        json_obj = json.loads(line)
                        results.append(json_obj)
                    except json.JSONDecodeError as e:
                        logger.warning(
                            f"JSON parsing error at line {line_num}: {str(e)}"
                        )

            # Return parsing results
            return json.dumps(
                {"format": "jsonl", "count": len(results), "data": results},
                ensure_ascii=False,
                cls=ComplexEncoder,
            )
        else:
            # Process standard JSON file
            with open(document_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            # Provide array length information if it's an array
            if isinstance(data, list):
                result = {
                    "format": "json",
                    "type": "array",
                    "count": len(data),
                    "data": data,
                }
            else:
                result = {
                    "format": "json",
                    "type": "object",
                    "keys": list(data.keys()) if isinstance(data, dict) else [],
                    "data": data,
                }

            return json.dumps(result, ensure_ascii=False, cls=ComplexEncoder)

    except json.JSONDecodeError as e:
        return handle_error(e, "JSON parsing")
    except Exception as e:
        return handle_error(e, "JSON file reading")


def mcpreadxml(
    document_path: str = Field(description="The input XML file path."),
) -> str:
    """Read and return content from XML file."""
    error = check_file_readable(document_path)
    if error:
        return json.dumps({"error": error})

    try:
        import xmltodict

        with open(document_path, "r", encoding="utf-8") as f:
            data = f.read()
        return json.dumps({"content": xmltodict.parse(data)}, ensure_ascii=False)
    except Exception as e:
        return handle_error(e, "XML file reading")


def mcpreadpdf(
    document_path: str = Field(description="The input PDF file path."),
) -> str:
    """Read and return content from PDF file."""
    error = check_file_readable(document_path)
    if error:
        return json.dumps({"error": error})

    try:
        from PyPDF2 import PdfReader

        with open(document_path, "rb") as f:
            reader = PdfReader(f)
            content = " ".join(page.extract_text() for page in reader.pages)
        return json.dumps({"content": content})
    except Exception as e:
        return handle_error(e, "PDF file reading")


def mcpreaddocx(
    document_path: str = Field(description="The input Word file path."),
) -> str:
    """Read and return content from Word file."""
    error = check_file_readable(document_path)
    if error:
        return json.dumps({"error": error})

    try:
        from docx2markdown._docx_to_markdown import docx_to_markdown

        file_name = os.path.basename(document_path)
        md_file_path = f"{file_name}.md"
        docx_to_markdown(document_path, md_file_path)
        with open(md_file_path, "r") as f:
            content = f.read()
        os.remove(md_file_path)
        return json.dumps({"content": content})
    except Exception as e:
        return handle_error(e, "Word file reading")


def mcpreadexcel(
    document_path: str = Field(description="The input Excel file path."),
) -> str:
    """Read and return content from Excel file."""
    error = check_file_readable(document_path)
    if error:
        return json.dumps({"error": error})

    try:
        import pandas as pd

        excel_data = {}

        with pd.ExcelFile(document_path) as xls:
            sheet_names = xls.sheet_names
            for sheet_name in sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                sheet_data = df.to_dict(orient="records")
                excel_data[sheet_name] = sheet_data

        return json.dumps(excel_data, ensure_ascii=False, cls=ComplexEncoder)
    except Exception as e:
        return handle_error(e, "Excel file reading")


def mcpreadpptx(
    document_path: str = Field(description="The input PowerPoint file path."),
) -> str:
    """Read and convert PowerPoint slides to base64 encoded images."""
    error = check_file_readable(document_path)
    if error:
        return json.dumps({"error": error})

    # Create temporary directory
    temp_dir = tempfile.mkdtemp()
    slides_data = []

    try:
        presentation = Presentation(document_path)
        total_slides = len(presentation.slides)

        if total_slides == 0:
            raise ValueError("PPT file does not contain any slides")

        # Process each slide
        for i, slide in enumerate(presentation.slides):
            # Set slide dimensions
            slide_width_px = 1920  # 16:9 ratio
            slide_height_px = 1080

            # Create blank image
            slide_img = Image.new("RGB", (slide_width_px, slide_height_px), "white")
            draw = ImageDraw.Draw(slide_img)
            font = ImageFont.load_default()

            # Draw slide number
            draw.text((20, 20), f"Slide {i+1}/{total_slides}", fill="black", font=font)

            # Process shapes in the slide
            for shape in slide.shapes:
                try:
                    # Process images
                    if hasattr(shape, "image") and shape.image:
                        image_stream = io.BytesIO(shape.image.blob)
                        img = Image.open(image_stream)
                        left = int(
                            shape.left * slide_width_px / presentation.slide_width
                        )
                        top = int(
                            shape.top * slide_height_px / presentation.slide_height
                        )
                        slide_img.paste(img, (left, top))

                    # Process text
                    elif hasattr(shape, "text") and shape.text:
                        text_left = int(
                            shape.left * slide_width_px / presentation.slide_width
                        )
                        text_top = int(
                            shape.top * slide_height_px / presentation.slide_height
                        )
                        draw.text(
                            (text_left, text_top),
                            shape.text,
                            fill="black",
                            font=font,
                        )

                except Exception as shape_error:
                    logger.warning(
                        f"Error processing shape in slide {i+1}: {str(shape_error)}"
                    )

            # Save slide image
            img_path = os.path.join(temp_dir, f"slide_{i+1}.jpg")
            slide_img.save(img_path, "JPEG")

            # Convert to base64
            base64_image = encode_image_from_file(img_path)
            slides_data.append(
                {
                    "slide_number": i + 1,
                    "image": f"data:image/jpeg;base64,{base64_image}",
                }
            )

    except Exception as e:
        return handle_error(e, "PowerPoint processing")
    finally:
        # Clean up temporary files
        try:
            for file in os.listdir(temp_dir):
                os.remove(os.path.join(temp_dir, file))
            os.rmdir(temp_dir)
        except Exception as cleanup_error:
            logger.warning(f"Error cleaning up temporary files: {str(cleanup_error)}")

    return json.dumps(slides_data, ensure_ascii=False)


def mcpreadsourcecode(
    document_path: str = Field(
        description="Source code file path, supports various programming languages."
    ),
) -> str:
    """Read and analyze source code file, return code content and basic structure information."""
    error = check_file_readable(document_path)
    if error:
        return json.dumps({"error": error})

    try:
        # Get file extension
        ext = os.path.splitext(document_path)[1].lower()

        # Read file content
        with open(document_path, "r", encoding="utf-8") as f:
            content = f.read()

        # Basic code information
        code_info = {
            "content": content,
            "file_type": ext,
            "line_count": len(content.splitlines()),
            "size_bytes": os.path.getsize(document_path),
            "last_modified": datetime.fromtimestamp(
                os.path.getmtime(document_path)
            ).strftime("%Y-%m-%d %H:%M:%S"),
        }

        # Special processing for different languages
        if ext in [".py"]:
            # Python file analysis
            import ast

            try:
                tree = ast.parse(content)
                code_info.update(
                    {
                        "classes": [
                            node.name
                            for node in ast.walk(tree)
                            if isinstance(node, ast.ClassDef)
                        ],
                        "functions": [
                            node.name
                            for node in ast.walk(tree)
                            if isinstance(node, ast.FunctionDef)
                        ],
                        "imports": [
                            node.names[0].name
                            for node in ast.walk(tree)
                            if isinstance(node, ast.Import)
                        ],
                    }
                )
            except Exception as e:
                logger.warning(f"Python AST analysis failed: {str(e)}")

        elif ext in [".js", ".ts"]:
            # JavaScript/TypeScript file
            import re

            code_info.update(
                {
                    "functions": re.findall(r"function\s+(\w+)", content),
                    "classes": re.findall(r"class\s+(\w+)", content),
                    "imports": re.findall(
                        r'import\s+.*?from\s+[\'"](.+?)[\'"]', content
                    ),
                }
            )

        elif ext in [".java"]:
            # Java file
            import re

            code_info.update(
                {
                    "package": re.findall(r"package\s+([\w.]+);", content),
                    "classes": re.findall(r"class\s+(\w+)", content),
                    "methods": re.findall(
                        r"(?:public|private|protected)\s+\w+\s+(\w+)\s*\(", content
                    ),
                }
            )

        elif ext in [".cpp", ".hpp", ".h", ".cc"]:
            # C++ file
            import re

            code_info.update(
                {
                    "classes": re.findall(r"class\s+(\w+)", content),
                    "functions": re.findall(
                        r"(?:[\w:~]+)\s+(\w+)\s*\([^)]*\)\s*(?:const)?\s*{", content
                    ),
                    "includes": re.findall(r'#include\s+[<"](.+?)[>"]', content),
                }
            )

        return json.dumps(code_info, ensure_ascii=False, cls=ComplexEncoder)

    except Exception as e:
        return handle_error(e, "Source code file reading")


# Main function
if __name__ == "__main__":
    run_mcp_server(
        "Document Server",
        funcs=[
            mcpreadtext,
            mcpreadjson,
            mcpreadxml,
            mcpreadpdf,
            mcpreaddocx,
            mcpreadexcel,
            mcpreadpptx,
            mcpreadsourcecode,
        ],
        port=4444,
    )
