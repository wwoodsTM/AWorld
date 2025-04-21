"""
Document MCP Server

This module provides MCP server functionality for document processing and analysis.
It handles various document formats including:
- Text files
- PDF documents
- Word documents (DOCX)
- Excel spreadsheets
- PowerPoint presentations
- JSON and XML files
- Source code files

Each document type has specialized processing functions that extract content,
structure, and metadata. The server focuses on local file processing with
appropriate validation and error handling.

Main functions:
- read_text: Reads plain text files
- read_pdf: Reads PDF files with optional image extraction
- read_docx: Reads Word documents
- read_excel: Reads Excel spreadsheets
- read_pptx: Reads PowerPoint presentations
- read_json: Reads and parses JSON/JSONL files
- read_xml: Reads and parses XML files
- read_source_code: Reads and analyzes source code files
"""

import ast
import io
import json
import os
import re
import tempfile
import traceback
from datetime import date, datetime
from typing import Any, Dict, List, Optional

import fitz
import html2text
import pandas as pd
import xmltodict
from bs4 import BeautifulSoup
from docx2markdown._docx_to_markdown import docx_to_markdown
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pydantic import BaseModel, Field
from PyPDF2 import PdfReader
from tabulate import tabulate

from aworld.logs.util import logger
from aworld.mcp_servers.abc.base import MCPServerBase, mcp
from aworld.mcp_servers.image_server import ImageServer
from aworld.mcp_servers.utils import parse_port, run_mcp_server
from aworld.utils import import_package, import_packages


# Define model classes for different document types
class TextDocument(BaseModel):
    """Model representing a text document"""

    content: str
    file_path: str
    file_name: str
    file_size: int
    last_modified: str


class HtmlDocument(BaseModel):
    """Model representing an HTML document"""

    content: str  # Extracted text content
    html_content: str  # Original HTML content
    file_path: str
    file_name: str
    file_size: int
    last_modified: str
    title: Optional[str] = None
    links: Optional[List[Dict[str, str]]] = None
    images: Optional[List[Dict[str, str]]] = None
    tables: Optional[List[str]] = None
    markdown: Optional[str] = None  # HTML converted to Markdown format


class JsonDocument(BaseModel):
    """Model representing a JSON document"""

    format: str  # "json" or "jsonl"
    type: Optional[str] = None  # "array" or "object" for standard JSON
    count: Optional[int] = None
    keys: Optional[List[str]] = None
    data: Any
    file_path: str
    file_name: str


class XmlDocument(BaseModel):
    """Model representing an XML document"""

    content: Dict
    file_path: str
    file_name: str


class PdfImage(BaseModel):
    """Model representing an image extracted from a PDF"""

    page: int
    format: str
    width: int
    height: int
    path: str


class PdfDocument(BaseModel):
    """Model representing a PDF document"""

    content: str
    file_path: str
    file_name: str
    page_count: int
    images: Optional[List[PdfImage]] = None
    image_count: Optional[int] = None
    image_dir: Optional[str] = None
    error: Optional[str] = None


class PdfResult(BaseModel):
    """Model representing results from processing multiple PDF documents"""

    total_files: int
    success_count: int
    failed_count: int
    results: List[PdfDocument]


class DocxDocument(BaseModel):
    """Model representing a Word document"""

    content: str
    file_path: str
    file_name: str


class ExcelSheet(BaseModel):
    """Model representing a sheet in an Excel file"""

    name: str
    data: List[Dict[str, Any]]
    markdown_table: str
    row_count: int
    column_count: int


class ExcelDocument(BaseModel):
    """Model representing an Excel document"""

    file_name: str
    file_path: str
    processed_path: Optional[str] = None
    file_type: str
    sheet_count: int
    sheet_names: List[str]
    sheets: List[ExcelSheet]
    success: bool = True
    error: Optional[str] = None


class ExcelResult(BaseModel):
    """Model representing results from processing multiple Excel documents"""

    total_files: int
    success_count: int
    failed_count: int
    results: List[ExcelDocument]


class PowerPointSlide(BaseModel):
    """Model representing a slide in a PowerPoint presentation"""

    slide_number: int
    image: str  # Base64 encoded image


class PowerPointDocument(BaseModel):
    """Model representing a PowerPoint document"""

    file_path: str
    file_name: str
    slide_count: int
    slides: List[PowerPointSlide]


class SourceCodeDocument(BaseModel):
    """Model representing a source code document"""

    content: str
    file_type: str
    file_path: str
    file_name: str
    line_count: int
    size_bytes: int
    last_modified: str
    classes: Optional[List[str]] = None
    functions: Optional[List[str]] = None
    imports: Optional[List[str]] = None
    package: Optional[List[str]] = None
    methods: Optional[List[str]] = None
    includes: Optional[List[str]] = None


class DocumentError(BaseModel):
    """Model representing an error in document processing"""

    error: str
    file_path: Optional[str] = None
    file_name: Optional[str] = None


class ComplexEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, datetime):
            return obj.strftime("%Y-%m-%d %H:%M:%S")
        elif isinstance(obj, date):
            return obj.strftime("%Y-%m-%d")
        else:
            return json.JSONEncoder.default(self, obj)


class DocumentServer(MCPServerBase):
    """
    Document Server class for processing and analyzing various document formats.

    This class provides methods for reading and processing different document types
    including text files, PDFs, Word documents, Excel spreadsheets, PowerPoint
    presentations, JSON/XML files, and source code files.
    """

    _instance = None

    def __new__(cls):
        """Implement singleton pattern"""
        if cls._instance is None:
            cls._instance = super(DocumentServer, cls).__new__(cls)
            cls._instance._init_server()
        return cls._instance

    def _init_server(self):
        """Initialize the document server with required packages"""
        import_package("cv2", install_name="opencv-python")
        import_package("fitz", install_name="PyMuPDF")
        import_packages(["xmltodict", "pandas", "docx2markdown", "PyPDF2"])
        import_package("bs4", install_name="beautifulsoup4")
        import_package("html2text", install_name="html2text")
        import_package("tabulate")
        logger.info("DocumentServer initialized")

    @classmethod
    def get_instance(cls):
        """Get the singleton instance of DocumentServer"""
        if cls._instance is None:
            return cls()
        return cls._instance

    @staticmethod
    def handle_error(
        e: Exception, error_type: str, file_path: Optional[str] = None
    ) -> str:
        """Unified error handling and return standard format error message"""
        error_msg = f"{error_type} error: {str(e)}"
        logger.error(traceback.format_exc())

        error = DocumentError(
            error=error_msg,
            file_path=file_path,
            file_name=os.path.basename(file_path) if file_path else None,
        )

        return error.model_dump_json()

    @staticmethod
    def check_file_readable(document_path: str) -> str:
        """Check if file exists and is readable, return error message or None"""
        if not os.path.exists(document_path):
            return f"File does not exist: {document_path}"
        if not os.access(document_path, os.R_OK):
            return f"File is not readable: {document_path}"
        return None

    @mcp
    @classmethod
    def read_text(
        cls,
        document_path: str = Field(description="The input local text file path."),
    ) -> str:
        """Read and return content from local text file. Cannot process https://URLs files."""
        # Handle Field objects if they're passed directly
        if hasattr(document_path, "default") and not isinstance(document_path, str):
            document_path = document_path.default

        error = cls.check_file_readable(document_path)
        if error:
            return DocumentError(error=error, file_path=document_path).model_dump_json()

        try:
            with open(document_path, "r", encoding="utf-8") as f:
                content = f.read()

            result = TextDocument(
                content=content,
                file_path=document_path,
                file_name=os.path.basename(document_path),
                file_size=os.path.getsize(document_path),
                last_modified=datetime.fromtimestamp(
                    os.path.getmtime(document_path)
                ).strftime("%Y-%m-%d %H:%M:%S"),
            )

            return result.model_dump_json()

        except Exception as e:
            return cls.handle_error(e, "Text file reading", document_path)

    @mcp
    @classmethod
    def read_json(
        cls,
        document_path: str = Field(description="Local path to JSON or JSONL file"),
        is_jsonl: bool = Field(
            default=False,
            description="Whether the file is in JSONL format (one JSON object per line)",
        ),
    ) -> str:
        """Read and parse JSON or JSONL file, return the parsed content. Cannot process https://URLs files."""
        # Handle Field objects if they're passed directly
        if hasattr(document_path, "default") and not isinstance(document_path, str):
            document_path = document_path.default

        if hasattr(is_jsonl, "default") and not isinstance(is_jsonl, bool):
            is_jsonl = is_jsonl.default

        error = cls.check_file_readable(document_path)
        if error:
            return DocumentError(error=error, file_path=document_path).model_dump_json()

        try:
            # Choose processing method based on file type
            if is_jsonl:
                # Process JSONL file (one JSON object per line)
                results = []
                with open(document_path, "r", encoding="utf-8") as f:
                    for line_num, line in enumerate(f, 1):
                        line = line.strip()
                        if not line:
                            continue
                        try:
                            json_obj = json.loads(line)
                            results.append(json_obj)
                        except json.JSONDecodeError as e:
                            logger.warning(
                                f"JSON parsing error at line {line_num}: {str(e)}"
                            )

                # Create result model
                result = JsonDocument(
                    format="jsonl",
                    count=len(results),
                    data=results,
                    file_path=document_path,
                    file_name=os.path.basename(document_path),
                )

            else:
                # Process standard JSON file
                with open(document_path, "r", encoding="utf-8") as f:
                    data = json.load(f)

                # Create result model based on data type
                if isinstance(data, list):
                    result = JsonDocument(
                        format="json",
                        type="array",
                        count=len(data),
                        data=data,
                        file_path=document_path,
                        file_name=os.path.basename(document_path),
                    )
                else:
                    result = JsonDocument(
                        format="json",
                        type="object",
                        keys=list(data.keys()) if isinstance(data, dict) else [],
                        data=data,
                        file_path=document_path,
                        file_name=os.path.basename(document_path),
                    )

            return result.model_dump_json()

        except json.JSONDecodeError as e:
            return cls.handle_error(e, "JSON parsing", document_path)
        except Exception as e:
            return cls.handle_error(e, "JSON file reading", document_path)

    @mcp
    @classmethod
    def read_xml(
        cls,
        document_path: str = Field(description="The local input XML file path."),
    ) -> str:
        """Read and return content from XML file. Cannot process https://URLs files."""
        # Handle Field objects if they're passed directly
        if hasattr(document_path, "default") and not isinstance(document_path, str):
            document_path = document_path.default

        error = cls.check_file_readable(document_path)
        if error:
            return DocumentError(error=error, file_path=document_path).model_dump_json()

        try:
            with open(document_path, "r", encoding="utf-8") as f:
                data = f.read()

            result = XmlDocument(
                content=xmltodict.parse(data),
                file_path=document_path,
                file_name=os.path.basename(document_path),
            )

            return result.model_dump_json()
        except Exception as e:
            return cls.handle_error(e, "XML file reading", document_path)

    @mcp
    @classmethod
    def read_pdf(
        cls,
        document_paths: List[str] = Field(
            description="The local input PDF file paths."
        ),
        extract_images: bool = Field(
            default=False,
            description="Whether to extract images from PDF (default: False)",
        ),
    ) -> str:
        """Read and return content from PDF file with optional image extraction. Cannot process https://URLs files."""
        # Handle Field objects if they're passed directly
        if hasattr(document_paths, "default") and not isinstance(document_paths, list):
            document_paths = document_paths.default

        if hasattr(extract_images, "default") and not isinstance(extract_images, bool):
            extract_images = extract_images.default

        try:
            results = []
            success_count = 0
            failed_count = 0

            for document_path in document_paths:
                error = cls.check_file_readable(document_path)
                if error:
                    results.append(
                        PdfDocument(
                            content="",
                            file_path=document_path,
                            file_name=os.path.basename(document_path),
                            page_count=0,
                            error=error,
                        )
                    )
                    failed_count += 1
                    continue

                try:
                    with open(document_path, "rb") as f:
                        reader = PdfReader(f)
                        content = " ".join(page.extract_text() for page in reader.pages)
                        page_count = len(reader.pages)

                        pdf_result = PdfDocument(
                            content=content,
                            file_path=document_path,
                            file_name=os.path.basename(document_path),
                            page_count=page_count,
                        )

                        # Extract images if requested
                        if extract_images:
                            images_data = []
                            # Use /tmp directory for storing images
                            output_dir = "/tmp/pdf_images"

                            # Create output directory if it doesn't exist
                            os.makedirs(output_dir, exist_ok=True)

                            # Generate a unique subfolder based on filename to avoid conflicts
                            pdf_name = os.path.splitext(
                                os.path.basename(document_path)
                            )[0]
                            timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
                            image_dir = os.path.join(
                                output_dir, f"{pdf_name}_{timestamp}"
                            )
                            os.makedirs(image_dir, exist_ok=True)

                            try:
                                # Open PDF with PyMuPDF
                                pdf_document = fitz.open(document_path)

                                # Iterate through each page
                                for page_index in range(len(pdf_document)):
                                    page = pdf_document[page_index]

                                    # Get image list
                                    image_list = page.get_images(full=True)

                                    # Process each image
                                    for img_index, img in enumerate(image_list):
                                        # Extract image information
                                        xref = img[0]
                                        base_image = pdf_document.extract_image(xref)
                                        image_bytes = base_image["image"]
                                        image_ext = base_image["ext"]

                                        # Save image to file in /tmp directory
                                        img_filename = f"pdf_image_p{page_index+1}_{img_index+1}.{image_ext}"
                                        img_path = os.path.join(image_dir, img_filename)

                                        with open(img_path, "wb") as img_file:
                                            img_file.write(image_bytes)
                                            logger.success(f"Image saved: {img_path}")

                                        # Get image dimensions
                                        with Image.open(img_path) as img:
                                            width, height = img.size

                                        # Add to results with file path instead of base64
                                        images_data.append(
                                            PdfImage(
                                                page=page_index + 1,
                                                format=image_ext,
                                                width=width,
                                                height=height,
                                                path=img_path,
                                            )
                                        )

                                pdf_result.images = images_data
                                pdf_result.image_count = len(images_data)
                                pdf_result.image_dir = image_dir

                            except Exception as img_error:
                                logger.error(
                                    f"Error extracting images: {str(img_error)}"
                                )
                                # Don't clean up on error so we can keep any successfully extracted images
                                pdf_result.error = str(img_error)

                    results.append(pdf_result)
                    success_count += 1

                except Exception as e:
                    results.append(
                        PdfDocument(
                            content="",
                            file_path=document_path,
                            file_name=os.path.basename(document_path),
                            page_count=0,
                            error=str(e),
                        )
                    )
                    failed_count += 1

            # Create final result
            pdf_result = PdfResult(
                total_files=len(document_paths),
                success_count=success_count,
                failed_count=failed_count,
                results=results,
            )

            return pdf_result.model_dump_json()

        except Exception as e:
            return cls.handle_error(e, "PDF file reading")

    @mcp
    @classmethod
    def read_docx(
        cls,
        document_path: str = Field(description="The local input Word file path."),
    ) -> str:
        """Read and return content from Word file. Cannot process https://URLs files."""
        # Handle Field objects if they're passed directly
        if hasattr(document_path, "default") and not isinstance(document_path, str):
            document_path = document_path.default

        error = cls.check_file_readable(document_path)
        if error:
            return DocumentError(error=error, file_path=document_path).model_dump_json()

        try:
            file_name = os.path.basename(document_path)
            md_file_path = f"{file_name}.md"
            docx_to_markdown(document_path, md_file_path)

            with open(md_file_path, "r") as f:
                content = f.read()

            os.remove(md_file_path)

            result = DocxDocument(
                content=content, file_path=document_path, file_name=file_name
            )

            return result.model_dump_json()
        except Exception as e:
            return cls.handle_error(e, "Word file reading", document_path)

    @mcp
    @classmethod
    def read_excel(
        cls,
        document_paths: List[str] = Field(
            description="List of local input Excel/CSV file paths."
        ),
        max_rows: int = Field(
            1000, description="Maximum number of rows to read per sheet (default: 1000)"
        ),
        convert_xls_to_xlsx: bool = Field(
            False,
            description="Whether to convert XLS files to XLSX format (default: False)",
        ),
    ) -> str:
        """Read multiple Excel/CSV files and convert sheets to Markdown tables. Cannot process https://URLs files."""
        # Handle Field objects if they're passed directly
        if hasattr(document_paths, "default") and not isinstance(document_paths, list):
            document_paths = document_paths.default

        if hasattr(max_rows, "default") and not isinstance(max_rows, int):
            max_rows = max_rows.default

        if hasattr(convert_xls_to_xlsx, "default") and not isinstance(
            convert_xls_to_xlsx, bool
        ):
            convert_xls_to_xlsx = convert_xls_to_xlsx.default

        try:
            # Import xls2xlsx package if conversion is requested
            if convert_xls_to_xlsx:
                import_package("xls2xlsx")
                from xls2xlsx import XLS2XLSX

            all_results = []
            temp_files = []  # Track temporary files for cleanup
            success_count = 0
            failed_count = 0

            # Process each file
            for document_path in document_paths:
                # Check if file exists and is readable
                error = cls.check_file_readable(document_path)
                if error:
                    all_results.append(
                        ExcelDocument(
                            file_name=os.path.basename(document_path),
                            file_path=document_path,
                            file_type="UNKNOWN",
                            sheet_count=0,
                            sheet_names=[],
                            sheets=[],
                            success=False,
                            error=error,
                        )
                    )
                    failed_count += 1
                    continue

                try:
                    # Check file extension
                    file_ext = os.path.splitext(document_path)[1].lower()

                    # Validate file type
                    if file_ext not in [".csv", ".xls", ".xlsx", ".xlsm"]:
                        error_msg = f"Unsupported file format: {file_ext}. Only CSV, XLS, XLSX, and XLSM formats are supported."
                        all_results.append(
                            ExcelDocument(
                                file_name=os.path.basename(document_path),
                                file_path=document_path,
                                file_type=file_ext.replace(".", "").upper(),
                                sheet_count=0,
                                sheet_names=[],
                                sheets=[],
                                success=False,
                                error=error_msg,
                            )
                        )
                        failed_count += 1
                        continue

                    # Convert XLS to XLSX if requested and file is XLS
                    processed_path = document_path
                    if convert_xls_to_xlsx and file_ext == ".xls":
                        try:
                            logger.info(f"Converting XLS to XLSX: {document_path}")
                            converter = XLS2XLSX(document_path)
                            # Create temp file with xlsx extension
                            xlsx_path = (
                                os.path.splitext(document_path)[0] + "_converted.xlsx"
                            )
                            converter.to_xlsx(xlsx_path)
                            processed_path = xlsx_path
                            temp_files.append(xlsx_path)  # Track for cleanup
                            logger.success(f"Converted XLS to XLSX: {xlsx_path}")
                        except Exception as conv_error:
                            logger.error(
                                f"XLS to XLSX conversion error: {str(conv_error)}"
                            )
                            # Continue with original file if conversion fails

                    excel_sheets = []
                    sheet_names = []

                    # Handle CSV files differently
                    if file_ext == ".csv":
                        # For CSV files, create a single sheet with the file name
                        sheet_name = os.path.basename(document_path).replace(".csv", "")
                        df = pd.read_csv(processed_path, nrows=max_rows)

                        # Create markdown table
                        markdown_table = "*Empty table*"
                        if not df.empty:
                            headers = df.columns.tolist()
                            table_data = df.values.tolist()
                            markdown_table = tabulate(
                                table_data, headers=headers, tablefmt="pipe"
                            )

                            if len(df) >= max_rows:
                                markdown_table += (
                                    f"\n\n*Note: Table truncated to {max_rows} rows*"
                                )

                        # Create sheet model
                        excel_sheets.append(
                            ExcelSheet(
                                name=sheet_name,
                                data=df.to_dict(orient="records"),
                                markdown_table=markdown_table,
                                row_count=len(df),
                                column_count=len(df.columns),
                            )
                        )

                        sheet_names = [sheet_name]

                    else:
                        # For Excel files, process all sheets
                        with pd.ExcelFile(processed_path) as xls:
                            sheet_names = xls.sheet_names

                            for sheet_name in sheet_names:
                                # Read Excel sheet into DataFrame with row limit
                                df = pd.read_excel(
                                    xls, sheet_name=sheet_name, nrows=max_rows
                                )

                                # Create markdown table
                                markdown_table = "*Empty table*"
                                if not df.empty:
                                    headers = df.columns.tolist()
                                    table_data = df.values.tolist()
                                    markdown_table = tabulate(
                                        table_data, headers=headers, tablefmt="pipe"
                                    )

                                    if len(df) >= max_rows:
                                        markdown_table += f"\n\n*Note: Table truncated to {max_rows} rows*"

                                # Create sheet model
                                excel_sheets.append(
                                    ExcelSheet(
                                        name=sheet_name,
                                        data=df.to_dict(orient="records"),
                                        markdown_table=markdown_table,
                                        row_count=len(df),
                                        column_count=len(df.columns),
                                    )
                                )

                    # Create result for this file
                    file_result = ExcelDocument(
                        file_name=os.path.basename(document_path),
                        file_path=document_path,
                        processed_path=(
                            processed_path if processed_path != document_path else None
                        ),
                        file_type=file_ext.replace(".", "").upper(),
                        sheet_count=len(sheet_names),
                        sheet_names=sheet_names,
                        sheets=excel_sheets,
                        success=True,
                    )

                    all_results.append(file_result)
                    success_count += 1

                except Exception as file_error:
                    # Handle errors for individual files
                    error_msg = str(file_error)
                    logger.error(f"File reading error for {document_path}: {error_msg}")
                    all_results.append(
                        ExcelDocument(
                            file_name=os.path.basename(document_path),
                            file_path=document_path,
                            file_type=os.path.splitext(document_path)[1]
                            .replace(".", "")
                            .upper(),
                            sheet_count=0,
                            sheet_names=[],
                            sheets=[],
                            success=False,
                            error=error_msg,
                        )
                    )
                    failed_count += 1

            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                        logger.info(f"Removed temporary file: {temp_file}")
                except Exception as cleanup_error:
                    logger.warning(
                        f"Error cleaning up temporary file {temp_file}: {str(cleanup_error)}"
                    )

            # Create final result
            excel_result = ExcelResult(
                total_files=len(document_paths),
                success_count=success_count,
                failed_count=failed_count,
                results=all_results,
            )

            return excel_result.model_dump_json()

        except Exception as e:
            return cls.handle_error(e, "Excel/CSV files processing")

    @mcp
    @classmethod
    def read_pptx(
        cls,
        document_path: str = Field(description="The local input PowerPoint file path."),
    ) -> str:
        """Read and convert PowerPoint slides to base64 encoded images. Cannot process https://URLs files."""
        # Handle Field objects if they're passed directly
        if hasattr(document_path, "default") and not isinstance(document_path, str):
            document_path = document_path.default

        error = cls.check_file_readable(document_path)
        if error:
            return DocumentError(error=error, file_path=document_path).model_dump_json()

        # Create temporary directory
        temp_dir = tempfile.mkdtemp()
        slides_data = []

        try:
            presentation = Presentation(document_path)
            total_slides = len(presentation.slides)

            if total_slides == 0:
                raise ValueError("PPT file does not contain any slides")

            # Process each slide
            for i, slide in enumerate(presentation.slides):
                # Set slide dimensions
                slide_width_px = 1920  # 16:9 ratio
                slide_height_px = 1080

                # Create blank image
                slide_img = Image.new("RGB", (slide_width_px, slide_height_px), "white")
                draw = ImageDraw.Draw(slide_img)
                font = ImageFont.load_default()

                # Draw slide number
                draw.text(
                    (20, 20), f"Slide {i+1}/{total_slides}", fill="black", font=font
                )

                # Process shapes in the slide
                for shape in slide.shapes:
                    try:
                        # Process images
                        if hasattr(shape, "image") and shape.image:
                            image_stream = io.BytesIO(shape.image.blob)
                            img = Image.open(image_stream)
                            left = int(
                                shape.left * slide_width_px / presentation.slide_width
                            )
                            top = int(
                                shape.top * slide_height_px / presentation.slide_height
                            )
                            slide_img.paste(img, (left, top))

                        # Process text
                        elif hasattr(shape, "text") and shape.text:
                            text_left = int(
                                shape.left * slide_width_px / presentation.slide_width
                            )
                            text_top = int(
                                shape.top * slide_height_px / presentation.slide_height
                            )
                            draw.text(
                                (text_left, text_top),
                                shape.text,
                                fill="black",
                                font=font,
                            )

                    except Exception as shape_error:
                        logger.warning(
                            f"Error processing shape in slide {i+1}: {str(shape_error)}"
                        )

                # Save slide image
                img_path = os.path.join(temp_dir, f"slide_{i+1}.jpg")
                slide_img.save(img_path, "JPEG")

                # Convert to base64 using ImageServer's class method
                base64_image = ImageServer.encode_images(img_path)
                slides_data.append(
                    PowerPointSlide(
                        slide_number=i + 1,
                        image=f"data:image/jpeg;base64,{base64_image}",
                    )
                )

            # Create result
            result = PowerPointDocument(
                file_path=document_path,
                file_name=os.path.basename(document_path),
                slide_count=total_slides,
                slides=slides_data,
            )

            return result.model_dump_json()

        except Exception as e:
            return cls.handle_error(e, "PowerPoint processing", document_path)
        finally:
            # Clean up temporary files
            try:
                for file in os.listdir(temp_dir):
                    os.remove(os.path.join(temp_dir, file))
                os.rmdir(temp_dir)
            except Exception as cleanup_error:
                logger.warning(
                    f"Error cleaning up temporary files: {str(cleanup_error)}"
                )

    @mcp
    @classmethod
    def read_source_code(
        cls,
        document_path: str = Field(
            description="Source code file path, supports various programming languages."
        ),
    ) -> str:
        """Read and analyze source code file, return code content and basic structure information. Cannot process https://URLs files."""
        # Handle Field objects if they're passed directly
        if hasattr(document_path, "default") and not isinstance(document_path, str):
            document_path = document_path.default

        error = cls.check_file_readable(document_path)
        if error:
            return DocumentError(error=error, file_path=document_path).model_dump_json()

        try:
            # Get file extension
            ext = os.path.splitext(document_path)[1].lower()

            # Read file content
            with open(document_path, "r", encoding="utf-8") as f:
                content = f.read()

            # Basic code information
            code_info = SourceCodeDocument(
                content=content,
                file_type=ext,
                file_path=document_path,
                file_name=os.path.basename(document_path),
                line_count=len(content.splitlines()),
                size_bytes=os.path.getsize(document_path),
                last_modified=datetime.fromtimestamp(
                    os.path.getmtime(document_path)
                ).strftime("%Y-%m-%d %H:%M:%S"),
            )

            # Special processing for different languages
            if ext in [".py"]:
                # Python file analysis
                try:
                    tree = ast.parse(content)
                    code_info.classes = [
                        node.name
                        for node in ast.walk(tree)
                        if isinstance(node, ast.ClassDef)
                    ]
                    code_info.functions = [
                        node.name
                        for node in ast.walk(tree)
                        if isinstance(node, ast.FunctionDef)
                    ]
                    code_info.imports = [
                        node.names[0].name
                        for node in ast.walk(tree)
                        if isinstance(node, ast.Import)
                    ]
                except Exception as e:
                    logger.warning(f"Python AST analysis failed: {str(e)}")

            elif ext in [".js", ".ts"]:
                # JavaScript/TypeScript file
                code_info.functions = re.findall(r"function\s+(\w+)", content)
                code_info.classes = re.findall(r"class\s+(\w+)", content)
                code_info.imports = re.findall(
                    r'import\s+.*?from\s+[\'"](.+?)[\'"]', content
                )

            elif ext in [".java"]:
                # Java file
                code_info.package = re.findall(r"package\s+([\w.]+);", content)
                code_info.classes = re.findall(r"class\s+(\w+)", content)
                code_info.methods = re.findall(
                    r"(?:public|private|protected)\s+\w+\s+(\w+)\s*\(", content
                )

            elif ext in [".cpp", ".hpp", ".h", ".cc"]:
                # C++ file
                code_info.classes = re.findall(r"class\s+(\w+)", content)
                code_info.functions = re.findall(
                    r"(?:[\w:~]+)\s+(\w+)\s*\([^)]*\)\s*(?:const)?\s*{", content
                )
                code_info.includes = re.findall(r'#include\s+[<"](.+?)[>"]', content)

            return code_info.model_dump_json()

        except Exception as e:
            return cls.handle_error(e, "Source code file reading")

    @classmethod
    def read_html_text(
        cls,
        document_path: str = Field(description="Local HTML file path or Web URL."),
        extract_links: bool = Field(
            default=True, description="Whether to extract link information"
        ),
        extract_images: bool = Field(
            default=True, description="Whether to extract image information"
        ),
        extract_tables: bool = Field(
            default=True, description="Whether to extract table information"
        ),
        convert_to_markdown: bool = Field(
            default=True, description="Whether to convert HTML to Markdown format"
        ),
    ) -> str:
        """Read HTML file and extract text content, optionally extract links, images, and table information, and convert to Markdown format."""
        # Handle Field objects if they're passed directly
        if hasattr(document_path, "default") and not isinstance(document_path, str):
            document_path = document_path.default

        if hasattr(extract_links, "default") and not isinstance(extract_links, bool):
            extract_links = extract_links.default

        if hasattr(extract_images, "default") and not isinstance(extract_images, bool):
            extract_images = extract_images.default

        if hasattr(extract_tables, "default") and not isinstance(extract_tables, bool):
            extract_tables = extract_tables.default

        if hasattr(convert_to_markdown, "default") and not isinstance(
            convert_to_markdown, bool
        ):
            convert_to_markdown = convert_to_markdown.default

        error = cls.check_file_readable(document_path)
        if error:
            return DocumentError(error=error, file_path=document_path).model_dump_json()

        try:
            # Read HTML file
            with open(document_path, "r", encoding="utf-8") as f:
                html_content = f.read()

            # Parse HTML using BeautifulSoup
            soup = BeautifulSoup(html_content, "html.parser")

            # Extract text content (remove script and style content)
            for script in soup(["script", "style"]):
                script.extract()
            text_content = soup.get_text(separator="\n", strip=True)

            # Extract title
            title = soup.title.string if soup.title else None

            # Initialize result object
            result = HtmlDocument(
                content=text_content,
                html_content=html_content,
                file_path=document_path,
                file_name=os.path.basename(document_path),
                file_size=os.path.getsize(document_path),
                last_modified=datetime.fromtimestamp(
                    os.path.getmtime(document_path)
                ).strftime("%Y-%m-%d %H:%M:%S"),
                title=title,
            )

            # Extract links
            if extract_links:
                links = []
                for link in soup.find_all("a"):
                    href = link.get("href")
                    text = link.get_text(strip=True)
                    if href:
                        links.append({"url": href, "text": text})
                result.links = links

            # Extract images
            if extract_images:
                images = []
                for img in soup.find_all("img"):
                    src = img.get("src")
                    alt = img.get("alt", "")
                    if src:
                        images.append({"src": src, "alt": alt})
                result.images = images

            # Extract tables
            if extract_tables:
                tables = []
                for table in soup.find_all("table"):
                    tables.append(str(table))
                result.tables = tables

            # Convert to Markdown
            if convert_to_markdown:
                h = html2text.HTML2Text()
                h.ignore_links = False
                h.ignore_images = False
                h.ignore_tables = False
                markdown_content = h.handle(html_content)
                result.markdown = markdown_content

            return result.model_dump_json()

        except Exception as e:
            return cls.handle_error(e, "HTML file reading", document_path)

    @mcp
    @classmethod
    def unzip_file(
        cls,
        zip_path: str = Field(description="Path to the ZIP file to extract"),
        extract_to: str = Field(
            default=None,
            description="Directory to extract files to. If None, extracts to the same directory as the ZIP file",
        ),
        password: str = Field(
            default=None, description="Password for encrypted ZIP files (if needed)"
        ),
    ) -> str:
        """
        Extract contents from a ZIP file to a specified directory.

        Args:
            zip_path: Path to the ZIP file to extract
            extract_to: Directory to extract files to (defaults to same directory as ZIP)
            password: Password for encrypted ZIP files (if needed)

        Returns:
            JSON string with extraction results
        """
        # Handle Field objects if they're passed directly
        if hasattr(zip_path, "default") and not isinstance(zip_path, str):
            zip_path = zip_path.default

        if (
            hasattr(extract_to, "default")
            and extract_to is not None
            and not isinstance(extract_to, str)
        ):
            extract_to = extract_to.default

        if (
            hasattr(password, "default")
            and password is not None
            and not isinstance(password, str)
        ):
            password = password.default

        error = cls.check_file_readable(zip_path)
        if error:
            return DocumentError(error=error, file_path=zip_path).model_dump_json()

        try:
            import zipfile

            # Determine extraction directory
            if not extract_to:
                extract_to = os.path.dirname(zip_path)

            # Create extraction directory if it doesn't exist
            os.makedirs(extract_to, exist_ok=True)

            # Track extraction results
            extracted_files = []
            failed_files = []

            with zipfile.ZipFile(zip_path, "r") as zip_ref:
                # Get list of all files in the archive
                file_list = zip_ref.namelist()

                # Extract each file
                for file in file_list:
                    try:
                        # Handle password if provided
                        if password:
                            zip_ref.extract(
                                file, path=extract_to, pwd=password.encode()
                            )
                        else:
                            zip_ref.extract(file, path=extract_to)

                        extracted_files.append(
                            {"name": file, "path": os.path.join(extract_to, file)}
                        )
                    except Exception as e:
                        failed_files.append({"name": file, "error": str(e)})

            # Prepare result
            result = {
                "success": True,
                "zip_file": zip_path,
                "extract_directory": extract_to,
                "total_files": len(file_list),
                "extracted_count": len(extracted_files),
                "failed_count": len(failed_files),
                "extracted_files": extracted_files,
                "failed_files": failed_files,
            }

            return json.dumps(result, cls=ComplexEncoder)

        except Exception as e:
            return cls.handle_error(e, "ZIP extraction", zip_path)


# Main function
if __name__ == "__main__":
    port = parse_port()

    document_server = DocumentServer.get_instance()
    logger.info("DocumentServer initialized and ready to handle requests")

    result = document_server.unzip_file(
        "/Users/arac/Desktop/gaia-benchmark/GAIA/2023/validation/9b54f9d9-35ee-4a14-b62f-d130ea00317f.zip",
        "/Users/arac/Desktop/",
    )
